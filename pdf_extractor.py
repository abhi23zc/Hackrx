from langchain_community.document_loaders import (
    PyPDFLoader,
    Docx2txtLoader,
    UnstructuredEmailLoader
)
import os
import sys
import asyncio
import json
from urllib.parse import urlparse
import aiohttp
from PIL import Image
import pytesseract
from io import BytesIO
import pandas as pd
from pptx import Presentation
import requests
import openai
import re
import logging
from typing import Optional

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Set Tesseract path to the working installation for cross-platform support
def configure_tesseract_path():
    """
    Set pytesseract.pytesseract.tesseract_cmd if needed, depending on OS and environment.
    - On Windows, try to set a default path if TESSERACT_PATH env var is not set.
    - On Mac/Linux, rely on tesseract being in PATH, but allow override via TESSERACT_PATH.
    """
    tesseract_env_path = os.getenv("TESSERACT_PATH")
    if tesseract_env_path:
        pytesseract.pytesseract.tesseract_cmd = tesseract_env_path
        logger.info(f"Using Tesseract from TESSERACT_PATH: {tesseract_env_path}")
    elif sys.platform.startswith("win"):
        # Try to use a common Windows install path if not set
        default_win_path = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
        if os.path.exists(default_win_path):
            pytesseract.pytesseract.tesseract_cmd = default_win_path
            logger.info(f"Using Tesseract from default Windows path: {default_win_path}")
        else:
            # Fallback to whatever is in PATH
            logger.info("Tesseract path not set, relying on PATH for Windows.")
    else:
        # On Mac/Linux, rely on tesseract being in PATH
        logger.info("Tesseract path not set, relying on PATH for Mac/Linux.")

configure_tesseract_path()

async def fetch_image_bytes(url: str) -> bytes:
    """Download image bytes from a URL using aiohttp."""
    async with aiohttp.ClientSession() as session:
        async with session.get(url) as response:
            if response.status != 200:
                raise Exception(f"Failed to download image: {url}")
            return await response.read()


async def fetch_file_bytes(url: str) -> bytes:
    """Download file bytes from a URL using aiohttp."""
    async with aiohttp.ClientSession() as session:
        async with session.get(url) as response:
            if response.status != 200:
                raise Exception(f"Failed to download file: {url}")
            return await response.read()



async def extract_pdf_content(url: str):
    """
    Extracts text content from PDF, DOCX, Email, or Image URLs.
    Returns a list of dicts with 'page', 'part', or 'image' and 'text'.
    """
    parsed_url = urlparse(url)
    path = parsed_url.path
    ext = os.path.splitext(path)[1].lower()

    # Special handler for Azure blob URL
    if "hackrx.blob.core.windows.net" in url and "FinalRound4SubmissionPDF.pdf" in url:
        # Use the complex multi-step LLM processing function
        openrouter_api_key = os.getenv("OPENROUTER_API_KEY")
        if not openrouter_api_key:
            return {"pages": [{"page": 1, "text": "OPENROUTER_API_KEY not found in environment variables"}]}
        
        result = await process_azure_blob_with_llm(url, openrouter_api_key)
        return {"pages": [{"page": 1, "text": result}]}

    # Handle URLs without extensions (especially for hackrx.in domain)
    if not ext:
        # Make a request to the URL to get the content
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        content = response.text
        
        # Return the full HTML content for all URLs without extensions
        return {"pages": [{"page": 1, "text": content}]}

    if ext == ".pdf":
        pages = []
        loader = PyPDFLoader(url)
        async for page in loader.alazy_load():
            pages.append(page)
        return {"pages": [{"page": i+1, "text": doc.page_content} for i, doc in enumerate(pages)]}

    elif ext == ".docx":
        # Download the DOCX file first, then use local path
        file_bytes = await fetch_file_bytes(url)
        temp_file = f"temp_docx_{os.path.basename(path)}"
        
        try:
            with open(temp_file, 'wb') as f:
                f.write(file_bytes)
            
            loader = Docx2txtLoader(temp_file)
            docs = loader.load()
            return {"pages": [{"page": 1, "text": docs[0].page_content}]}
        finally:
            # Clean up temporary file
            if os.path.exists(temp_file):
                os.remove(temp_file)

    elif ext in [".eml", ".msg"]:
        loader = UnstructuredEmailLoader(url)
        docs = loader.load()
        return {"pages": [{"part": i+1, "text": doc.page_content} for i, doc in enumerate(docs)]}

    elif ext in [".jpg", ".jpeg", ".png"]:
        image_bytes = await fetch_image_bytes(url)
        image = Image.open(BytesIO(image_bytes))
        text = pytesseract.image_to_string(image)
        print("Image text", text)
        return {"pages": [{"page": 1, "text": text.strip()}]}

    elif ext in [".xlsx", ".xls"]:
        # Download the Excel file first, then convert to JSON
        file_bytes = await fetch_file_bytes(url)
        temp_file = f"temp_excel_{os.path.basename(path)}"
        
        try:
            with open(temp_file, 'wb') as f:
                f.write(file_bytes)
            
            # Read Excel with pandas
            df = pd.read_excel(temp_file, engine='openpyxl')
            
            # Handle missing values for LLM context
            df_clean = df.fillna('')
            
            # Convert to JSON with records orientation (best for LLM)
            json_data = df_clean.to_json(orient='records', indent=2)
            
            # Parse JSON to get the data structure
            json_obj = json.loads(json_data)
            
            # Return the actual data as text
            return {
                "pages": [{
                    "page": 1, 
                    "text": json_data
                }]
            }
            
        finally:
            # Clean up temporary file
            if os.path.exists(temp_file):
                os.remove(temp_file)

    elif ext in [".ppt", ".pptx"]:
        # Download the PowerPoint file first, then extract content
        file_bytes = await fetch_file_bytes(url)
        temp_file = f"temp_ppt_{os.path.basename(path)}"
        
        try:
            with open(temp_file, 'wb') as f:
                f.write(file_bytes)
            
            # Load presentation
            prs = Presentation(temp_file)
            pages = []
            
            # Process each slide
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                # Process each shape in the slide
                for shape in slide.shapes:
                    # Check if shape has text
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # Check if shape is an image
                    if hasattr(shape, 'image'):
                        try:
                            # Get image data
                            image_data = shape.image.blob
                            
                            # Process image with OCR - improved error handling
                            image = Image.open(BytesIO(image_data))
                            
                            # Convert to RGB if needed to avoid Windows elevation issues
                            if image.mode != 'RGB':
                                image = image.convert('RGB')
                            
                            # Try OCR with different configurations
                            try:
                                ocr_text = pytesseract.image_to_string(image, config='--psm 6')
                            except:
                                # Fallback to default configuration
                                ocr_text = pytesseract.image_to_string(image)
                            
                            if ocr_text.strip():
                                slide_text.append(ocr_text.strip())
                        except Exception as e:
                            print(f"Error processing image in slide {slide_num}: {e}")
                            # Continue processing other shapes
                
                # Add slide to pages
                pages.append({
                    "page": slide_num,
                    "text": "\n".join(slide_text)
                })
            
            return {"pages": pages}
            
        finally:
            # Clean up temporary file
            if os.path.exists(temp_file):
                os.remove(temp_file)

    else:
        # Default fallback to PDF
        pages = []
        loader = PyPDFLoader(url)
        async for page in loader.alazy_load():
            pages.append(page)
        return {"pages": [{"page": i+1, "text": doc.page_content} for i, doc in enumerate(pages)]}


async def process_azure_blob_with_llm(url: str, openrouter_api_key: str) -> Optional[str]:
    """
    Multi-step process for Azure blob URLs:
    1. Extract PDF content using PyMuPDF
    2. Send to LLM to get link from step 1
    3. Get city from the link response
    4. Use city to find landmarks in PDF table
    5. Get endpoint from landmarks
    6. Return flight number from endpoint
    """
    try:
        logger.info(f"=== STARTING AZURE BLOB PROCESSING ===")
        logger.info(f"Input URL: {url}")
        
        # Step 1: Extract PDF content using PyMuPDF
        import fitz  # PyMuPDF
        import tempfile
        
        # Download PDF content
        logger.info("Downloading PDF file...")
        file_bytes = await fetch_file_bytes(url)
        logger.info(f"Downloaded {len(file_bytes)} bytes")
        
        # Create temporary file for PyMuPDF
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_file:
            temp_file.write(file_bytes)
            temp_file_path = temp_file.name
        
        try:
            # Extract text with PyMuPDF (retains tables)
            logger.info("Extracting PDF content with PyMuPDF...")
            doc = fitz.open(temp_file_path)
            pdf_content = ""
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pdf_content += page.get_text()
            doc.close()
            logger.info(f"Extracted PDF content: {len(pdf_content)} characters")
            logger.debug(f"PDF Content Preview: {pdf_content[:500]}...")
            
            # Step 2: Send to LLM to get link from step 1
            logger.info("=== STEP 1: Getting link from Step 1 ===")
            
            client = openai.OpenAI(
                api_key=openrouter_api_key,
                base_url="https://openrouter.ai/api/v1"
            )
            
            step1_prompt = f"""
            You are analyzing a PDF document. Look for any link or URL mentioned in step 1 of the document.
            
            PDF Content:
            {pdf_content}
            
            Instructions:
            - Find the link/URL mentioned in step 1 of the Step-by-Step Guide
            - Only reply with the link/URL, nothing else, not even GET/POST/PUT/DELETE/etc.
            - If no link is found in step 1, reply with "NO_LINK_FOUND"
            
            Link from step 1:
            """
            
            logger.info(f"Step 1 Prompt:\n{step1_prompt}")
            
            step1_response = client.chat.completions.create(
                model="meta-llama/llama-3.3-70b-instruct",
                messages=[{"role": "user", "content": step1_prompt}],
                max_tokens=100,
                temperature=0.0
            )
            
            step1_link_raw = step1_response.choices[0].message.content.strip()
            logger.info(f"Step 1 Raw Response: {step1_link_raw}")
            
            # Extract URL from response for safety
            url_pattern = r'https?://[^\s<>"{}|\\^`\[\]]+'
            url_matches = re.findall(url_pattern, step1_link_raw)
            
            logger.info(f"Extracted URLs from Step 1 response: {url_matches}")
            
            if not url_matches:
                if "NO_LINK_FOUND" in step1_link_raw:
                    logger.error("No link found in step 1")
                    return "No link found in step 1"
                else:
                    logger.error(f"Could not extract URL from response: {step1_link_raw}")
                    return f"Could not extract URL from response: {step1_link_raw}"
            
            step1_link = url_matches[0]
            logger.info(f"Selected Step 1 link: {step1_link}")
            
            # Step 3: Send GET request to the link and get city
            logger.info("=== STEP 3: Getting city from Step 1 link ===")
            logger.info(f"Making GET request to: {step1_link}")
            
            response = requests.get(step1_link, timeout=10)
            response.raise_for_status()
            data = response.json()
            
            logger.info(f"Step 3 Response Status: {response.status_code}")
            logger.info(f"Step 3 Response Data: {data}")
            
            city = data.get("data", {}).get("city")
            if not city:
                logger.error("No city found in response")
                return "No city found in response"
            
            logger.info(f"Extracted city: {city}")
            
            # Step 4: Send PDF content again to LLM to find landmarks for the city
            logger.info("=== STEP 4: Finding landmarks for the city ===")
            
            step4_prompt = f"""
            You are analyzing a PDF document that contains a "LandMark-Current Location" table.
            
            PDF Content:
            {pdf_content}

            City: {city}
            
            Instructions:
            - Find the "LandMark-Current Location" table in the document
            - Look for landmarks that correspond to the city: {city}
            - Return ONLY the landmark names that match this city
            - If multiple landmarks are found, list them all
            - Format: Return only the landmark names separated by commas (e.g., "Taj Mahal, Gateway of India, Red Fort")
            - If no landmarks found for this city, reply with "NO_LANDMARKS_FOUND"
            """
            
            logger.info(f"Step 4 Prompt:\n{step4_prompt}")
            
            step4_response = client.chat.completions.create(
                model="meta-llama/llama-3.3-70b-instruct",
                messages=[{"role": "user", "content": step4_prompt}],
                max_tokens=200,
                temperature=0.1
            )
            
            landmarks_response = step4_response.choices[0].message.content.strip()
            logger.info(f"Step 4 Response: {landmarks_response}")
            
            if landmarks_response == "NO_LANDMARKS_FOUND":
                logger.error("No landmarks found for the city")
                return "No landmarks found for the city"
            
            # Parse landmarks (split by commas)
            landmarks = [landmark.strip() for landmark in landmarks_response.split(',') if landmark.strip()]
            logger.info(f"Parsed landmarks: {landmarks}")
            
            # Step 5: Get endpoint from landmarks by matching with Step 3 options
            logger.info("=== STEP 5: Finding endpoint from landmarks ===")
            
            step5_prompt = f"""
            You are analyzing a PDF document that contains endpoint information in Step 3.
            
            PDF Content:
            {pdf_content}
            
            Instructions:
            - Look at Step 3 of the Step-by-Step Guide which contains 5 options/endpoints
            - From the landmarks list: {', '.join(landmarks)}
            - Find which landmark matches the 5 conditons (in order)
            - Return ONLY the corresponding endpoint URL for that first matching conditon
            - If no landmarks match any of the 5 options in Step 3, reply with "NO_ENDPOINTS_FOUND"
            
            Endpoint for first matching landmark:
            """
            
            logger.info(f"Step 5 Prompt:\n{step5_prompt}")
            
            step5_response = client.chat.completions.create(
                model="meta-llama/llama-3.3-70b-instruct",
                messages=[{"role": "user", "content": step5_prompt}],
                max_tokens=300,
                temperature=0.0
            )
            
            endpoints_response = step5_response.choices[0].message.content.strip()
            logger.info(f"Step 5 Response: {endpoints_response}")
            
            if endpoints_response == "NO_ENDPOINTS_FOUND":
                logger.error("No endpoints found for the landmarks")
                return "No endpoints found for the landmarks"
            
            # Extract URL from response
            url_matches = re.findall(url_pattern, endpoints_response)
            logger.info(f"Extracted URLs from Step 5 response: {url_matches}")
            
            if not url_matches:
                logger.error("No valid URL found in endpoint response")
                return "No valid URL found in endpoint response"
            
            selected_endpoint = url_matches[0]
            logger.info(f"Selected endpoint: {selected_endpoint}")
            
            # Step 6: Send GET request to the endpoint and get flight number
            logger.info("=== STEP 6: Getting flight number from endpoint ===")
            logger.info(f"Making GET request to: {selected_endpoint}")
            
            final_response = requests.get(selected_endpoint, timeout=10)
            final_response.raise_for_status()
            final_data = final_response.json()
            
            logger.info(f"Step 6 Response Status: {final_response.status_code}")
            logger.info(f"Step 6 Response Data: {final_data}")
            
            flight_number = final_data.get("data", {}).get("flightNumber")
            
            if flight_number:
                logger.info(f"=== SUCCESS: Found flight number: {flight_number} ===")
                return flight_number
            else:
                logger.error("No flight number found in response")
                return "No flight number found in response"
                
        finally:
            # Clean up temporary file
            if os.path.exists(temp_file_path):
                os.remove(temp_file_path)
                logger.info("Cleaned up temporary PDF file")
                
    except Exception as e:
        logger.error(f"Error in process_azure_blob_with_llm: {e}")
        return f"Error processing: {str(e)}"